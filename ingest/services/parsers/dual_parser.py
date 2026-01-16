from pathlib import Path
import os, io, json
import logging
from typing import List, Dict, Any, Union
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
from bs4 import BeautifulSoup
from pptx import Presentation
from google import genai
from django.conf import settings
from retry import retry
import hashlib

# Configure logging
logger = logging.getLogger(__name__)

client = genai.Client(api_key=settings.GOOGLE_API_KEY)


# === TEXTUAL PATH ===
def extract_textual_content(file_path: str) -> List[Dict[str, Union[int, str]]]:
    logger.info(f"Extracting textual content from: {file_path}")
    text_data: List[Dict[str, Union[int, str]]] = []

    try:
        if file_path.endswith(".pdf"):
            doc = fitz.open(file_path)
            for i, page in enumerate(doc, 1):
                text = page.get_text("text").strip()
                if text:
                    text_data.append({"page": i, "text": text})

        elif file_path.endswith(".pptx"):
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides, 1):
                slide_text = " ".join(
                    shape.text for shape in slide.shapes if hasattr(shape, "text")
                )
                text_data.append({"slide": i, "text": slide_text})

        elif file_path.endswith(".html"):
            with open(file_path, "r", encoding="utf-8") as f:
                soup = BeautifulSoup(f, "html.parser")
            body_text = soup.get_text(separator=" ", strip=True)
            text_data.append({"section": "html", "text": body_text})

        logger.info(f"Extracted {len(text_data)} textual segments.")

    except Exception as e:
        logger.error(f"Error extracting text: {e}", exc_info=True)

    return text_data


# === VISUAL PATH ===
def extract_visual_content(file_path: str) -> List[Dict[str, Union[int, str]]]:
    logger.info(f"Extracting visual content from: {file_path}")
    visual_text: List[Dict[str, Union[int, str]]] = []

    try:
        images_to_process = []

        if file_path.endswith(".pdf"):
            doc = fitz.open(file_path)
            for i, page in enumerate(doc, 1):
                pix = page.get_pixmap(dpi=200)
                img_data = pix.tobytes("png")
                # Convert to PIL Image
                img = Image.open(io.BytesIO(img_data))
                images_to_process.append({"index": i, "image": img, "type": "page"})

        elif file_path.endswith(".pptx"):
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # picture
                        image = shape.image
                        image_bytes = io.BytesIO(image.blob)
                        img = Image.open(image_bytes)
                        images_to_process.append(
                            {"index": i, "image": img, "type": "slide"}
                        )

        # Process with Tesseract (Local OCR)
        for item in images_to_process:
            try:
                # Tesseract OCR
                ocr_text = pytesseract.image_to_string(item["image"])

                if ocr_text.strip():
                    visual_text.append(
                        {item["type"]: item["index"], "ocr_text": ocr_text.strip()}
                    )

            except Exception as inner_e:
                logger.warning(
                    f"Failed to process image {item['index']} with Tesseract: {inner_e}"
                )

        logger.info(f"Extracted {len(visual_text)} visual segments using Tesseract.")

    except Exception as e:
        logger.error(f"Error extracting visual content: {e}", exc_info=True)

    return visual_text


# === HIERARCHY REPAIR ===
@retry(tries=3, delay=1)
def fix_procedural_nesting(
    node: dict[str, Any] | list[Any],
) -> dict[str, Any] | list[Any]:
    """
    Repairs procedural step hierarchies and removes duplicates.
    Ensures only one top-level step1 per procedure, and subsequent steps
    nest inside one another (step2 → step3 → ...).
    """
    if not isinstance(node, dict):
        return node

    # Ensure a children array
    if "children" not in node or not isinstance(node["children"], list):
        node["children"] = []

    # Recursively fix all nested children first
    for child in node["children"]:
        fix_procedural_nesting(child)

    # Gather all procedural step nodes at this level
    step_nodes = [
        c
        for c in node["children"]
        if isinstance(c, dict) and "-step" in c.get("id", "")
    ]
    if not step_nodes:
        return node

    # Sort by step number (handles step10 safely)
    def step_num(n: dict[str, Any]) -> int:
        try:
            return int(n["id"].split("-step")[-1])
        except Exception:
            return 999

    step_nodes.sort(key=step_num)

    # Build a clean chain: step1 → step2 → step3 → ...
    for i in range(len(step_nodes) - 1):
        curr, nxt = step_nodes[i], step_nodes[i + 1]
        curr.setdefault("children", [])
        # Prevent duplicate reattachment
        if nxt["id"] not in [ch.get("id") for ch in curr["children"]]:
            curr["children"].append(nxt)

    # Only keep the first step at this level
    first_step = step_nodes[0]
    # Filter out all step duplicates from top-level children
    non_step_nodes = [c for c in node["children"] if c not in step_nodes]
    node["children"] = non_step_nodes + [first_step]

    return node


# === MERGE + GEMINI CALL ===
def parse_dualpath(file_path: str) -> str:
    logger.info(f"Starting dual-path parsing for: {file_path}")

    textual = extract_textual_content(file_path)
    visual = extract_visual_content(file_path)
    merged = {"textual": textual, "visual": visual}

    logger.info(f"Merged {len(textual)} textual and {len(visual)} visual segments")

    prompt = """
    You are an expert educational data modeler and AI architect specializing in automated knowledge graph construction for multi-agent learning systems. Your goal is to read and interpret any instructional or challenge-based document (e.g., Capture-the-Flag tutorials, lab guides, programming assignments, or course handouts) and automatically construct a unified Conceptual Knowledge Graph (KG) representing the learning structure contained within it.
    The output must be a single JSON object following the provided schema. Each concept node corresponds to a distinct atomic learning concept and internally embeds its procedural and assessment details as metadata.

    Generation Rules:
    Concept Nodes Only -
    The KG consists entirely of concept nodes (Cxx), each encapsulating its 1) Conceptual metadata (definition, Bloom level, difficulty, misconceptions, etc.), 2) Procedural metadata (steps, code snippets, etc.), and 3) Assessment metadata (questions to test student understanding, evaluation criteria, etc.)


    Procedural and Assessment Integration:
    There are no separate graphs for “procedural” or “assessment.”
    All procedures and evaluations appear as embedded children inside their corresponding concept. In this way, on clicking the conceptual nodes, their corresponding procedural and assessment nodes will be revealed for a more intuitive and usable UI.


    Automatic Coverage:
    Detect all challenge or module-level sections automatically (e.g., headings like “Challenge 01,” “Task 2,” or “Exercise 5”). Each distinct challenge should yield one or more concept nodes. Ensure that every challenge or instructional section in the document is represented through at least one concept node.


    Relations:
    Use semantic links between concepts such as: "PREREQUISITE_FOR", "DEPENDS_ON", "EXTENDS_TO", "GENERALIZES_TO", "USES_INIT_FROM", "IMPLEMENTED_BY", "ENABLES".


    These remain metadata (not graph edges).


    Pedagogical Encoding - 
    Bloom levels:


    Conceptual: Remember / Understand


    Procedural: Apply / Analyze


    Assessment: Evaluate / Create


    Visibility: ["supervisor_agent", "instructor"]


    Validation status: "pending", "verified", or "rejected"


    Confidence & relevance: values between 0.0 – 1.0


    Hinting Rules:
    Provide progressive scaffolding (general → specific → structural). NEVER REVEAL COMPLETE SOLUTIONS OR FULL CODE.
    Coverage Directive
    For any given input document:
    Identify all unique challenges or topics.


    For each, extract the main concept(s), associated procedures, and assessment components.


    The KG must comprehensively cover all detected sections, without having to explicitly list them (e.g., 00–10).


    Produce a single JSON object following this structure. It should work for any PDF and automatically adjust the number and naming of concept nodes based on the detected sections or challenges. The output of the JSON should STRICTLY adhere to this example structure strictly (the example below only has 2 conceptual nodes and covers 2 concepts. It could be more or less than this, as per needed). KEEP IN MIND THE NESTED STRUCTURE, IT IS VERY VERY IMPORTANT - All steps in procedural should be children of its previous step (e.g., "P01-step3" is a children of "P01-step2", which is a children of "P01-step1"). Refer to the example structure below and DO NOT MISS THIS:
    {
    "id": "CTF_KG",
    "name": "Central node",
    "children": [
        {
        "id": "C01",
        "name": "Angr Project Initialization",
        "label": "Concept 1",
        "definition": "Understanding how to initialize an Angr project and create an initial symbolic state for execution.",
        "difficulty": "beginner",
        "bloom_level": "Understand",
        "prerequisites": [],
        "misconceptions": ["Assuming angr loads libraries automatically", "Using full_init_state instead of entry_state"],
        "visibility": ["supervisor_agent", "instructor"],
        "validation_status": "verified",
        "confidence": 0.93,
        "relevance_score": 0.82,
        "source": "CTF_copy.pdf [page 5]",
        "learning_objective": "Set up an angr project and create a simulation state.",
        "connections": [
            {"to": "C02", "relation": "PREREQUISITE_FOR"}
        ],
        "children": [
            {
            "id": "P01",
            "name": "Procedural steps: Angr Initialization",
            "label": "Procedural Steps",
            "description": "Procedure for initializing an angr project and creating a symbolic state.",
            "difficulty": "beginner",
            "bloom_level": "Apply",
            "common_errors": ["Wrong binary path", "Using wrong initialization state"],
            "success_criteria": ["Project loads correctly", "Initial symbolic state created"],
            "error_patterns": ["FileNotFoundError", "AttributeError: 'Project' object has no attribute"],
            "progress_metric": {"completed": false, "percent_done": 0},
            "children": [
                {
                "id": "P01-step1",
                "name": "Step 1: Import angr and load binary",
                "label": "Procedural Step 1",
                "code_snippet": "import angr\nproject = angr.Project('challenge_binary', auto_load_libs=False)",
                "hint": "Ensure angr is installed; disable auto_load_libs to save memory.",
                "children": [
                    {
                    "id": "P01-step2",
                    "name": "Step 2: Create initial symbolic state",
                    "code_snippet": "state = project.factory.entry_state()",
                    "hint": "entry_state() begins at the binary entry point."
                    }
                ]
                }
            ]
            },
            {
            "id": "A01",
            "name": "Assessment Guide: Angr Initialization",
            "label": "Assessment 1",
            "linked_challenges": ["00_angr_find"],
            "objectives": [
                "Load binary in angr",
                "Initialize a correct symbolic state"
            ],
            "question_prompts": [
                {"question": "Which function creates a symbolic start state?"},
                {"question": "What happens if you use full_init_state()?"}
            ],
            "evaluation_criteria": [
                "Binary loads without crash",
                "State includes symbolic stdin"
            ],
            "bloom_level": "Evaluate",
            "difficulty": "beginner"
            }
        ]
        },
        {
        "id": "C02",
        "name": "Find and Avoid Algorithm",
        "definition": "Technique for identifying exploration targets in disassembly.",
        "label": "Concept 2",
        "difficulty": "intermediate",
        "bloom_level": "Apply",
        "prerequisites": ["C01"],
        "misconceptions": ["Find/avoid works automatically", "Using string address instead of code address"],
        "visibility": ["supervisor_agent", "instructor"],
        "validation_status": "verified",
        "confidence": 0.91,
        "relevance_score": 0.88,
        "source": "CTF_copy.pdf [page 9]",
        "learning_objective": "Use disassembly evidence to locate and configure find/avoid targets.",
        "connections": [
            {"to": "C01", "relation": "DEPENDS_ON"},
            {"to": "C03", "relation": "EXTENDS_TO"}
        ],
        "children": [
            {
            "id": "P02",
            "name": "Procedural Steps: Find/Avoid Algorithm",
            "label": "Procedural Steps",
            "description": "Procedure for extracting disassembly targets and configuring angr exploration.",
            "difficulty": "intermediate",
            "bloom_level": "Apply",
            "common_errors": [
                "Choosing unreachable addresses",
                "Incorrect explore() syntax"
            ],
            "success_criteria": [
                "Correct find/avoid addresses chosen",
                "Program reaches success function"
            ],
            "error_patterns": [
                "angr.SimManager Error",
                "InvalidAddressException"
            ],
            "progress_metric": {"completed": false, "percent_done": 0},
            "children": [
                {
                "id": "P02-step1",
                "name": "Step 1: Search for success message in disassembly",
                "label": "Procedural Step 1",
                "code_snippet": "objdump -d binary | grep 'puts'",
                "hint": "Look for string references near 'Success'.",
                "children": [
                    {
                    "id": "P02-step2",
                    "name": "Step 2: Find call site printing success message",
                    "label": "Procedural Step 2",
                    "code_snippet": "grep -n 'puts' binary.asm",
                    "hint": "Locate printf or puts thunks.",
                    "children": [
                        {
                        "id": "P02-step3",
                        "name": "Step 3: Record relevant instruction addresses",
                        "label": "Procedural Step 3",
                        "hint": "Mark addresses for find/avoid targets.",
                        "children": [
                            {
                            "id": "P02-step4",
                            "name": "Step 4: Configure find/avoid in explore()",
                            "label": "Procedural Step 4",
                            "code_snippet": "simgr.explore(find=success_addr, avoid=fail_addr)",
                            "hint": "Ensure both parameters are defined."
                            }
                        ]
                        }
                    ]
                    }
                ]
                }
            ]
            },
            {
            "id": "A02",
            "name": "Assessment Guide: Find/Avoid Algorithm",
            "label": "Assessment 2",
            "linked_challenges": ["02_angr_find_condition"],
            "objectives": [
                "Apply find/avoid correctly",
                "Link disassembly evidence to execution flow"
            ],
            "question_prompts": [
                {"question": "Why define both find and avoid addresses?"},
                {"question": "What happens if the avoid address is missed?"}
            ],
            "evaluation_criteria": [
                "Reaches correct success path",
                "Avoids failure outputs"
            ],
            "bloom_level": "Evaluate",
            "difficulty": "intermediate"
            }
            ]
          }
       ]
   }


    The following extracted content contains textual and OCR segments from the source material.
    Use it to populate the JSON fields accurately.
    """

    prompt += (
        "\nThe following extracted content contains textual and OCR segments from the source material. Use it to populate the JSON fields accurately:\n"
        + json.dumps(merged, indent=2)
    )

    logger.info("Sending prompt to Gemini models...")
    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash", contents=prompt
        )
        logger.info("Received response from Gemini.")
    except Exception as e:
        logger.error(f"Gemini API Error: {e}")
        raise e

    clean_json = response.text.strip()

    if clean_json.startswith("```"):
        clean_json = clean_json.replace("```json", "").replace("```", "").strip()

    try:
        parsed = json.loads(clean_json)

        base_name = os.path.splitext(os.path.basename(file_path))[0]

        def update_source_links(node):
            if isinstance(node, dict):
                # If this node has a "source" field like "CTF_copy.pdf [page 5]"
                if "source" in node and isinstance(node["source"], str):
                    node["source"] = f"/static/uploads/{base_name}.pdf"
                # Recurse through children
                for child in node.get("children", []):
                    update_source_links(child)
            elif isinstance(node, list):
                for n in node:
                    update_source_links(n)

        update_source_links(parsed)

        parsed = fix_procedural_nesting(parsed)
        content = json.dumps(parsed, ensure_ascii=False, indent=2)
        hash = hashlib.md5(content.encode(), usedforsecurity=False).hexdigest()
        output_path = Path(f"data/{hash}.json")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(content, encoding="utf-8")

        logger.info(f"Cleaned and saved valid JSON graph → {output_path}")

    except Exception as e:
        # Fallback: still save raw response so you can debug later
        logger.error(f"Gemini output not valid JSON ({e}). Saved raw text instead.")
        # with open(output_path, "w", encoding="utf-8") as f:
        #     f.write(response.text)
        # Raise error to fail the task if JSON invalid is critical
        raise ValueError("Model generated invalid JSON") from e

    return output_path
